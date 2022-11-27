from typing import Optional, List
from pathlib import Path
from itertools import chain
from pptx.enum.text import MSO_AUTO_SIZE
import wikipediaapi
import pptx

# Big thanks, I couldn't figure it out myself
# https://stackoverflow.com/a/12472564

def flatten(S):
    if S == []:
        return S
    if isinstance(S[0], list):
        return flatten(S[0]) + flatten(S[1:])
    return S[:1] + flatten(S[1:])

class Generator:
    def __init__(self, topics: List[str], lang: str = "pl", author:str = "Author", title:str = "Title", path: Path=Path("output.pptx"),max_letters_per_page: int = 500, split_on: str = "\n"):
        self.lang = lang
        self.path = path
        self.title = title
        self.author = author
        self.max_letters_per_page = max_letters_per_page
        self.split_on = split_on
        self.contents = []
        
        self.wikipedia = wikipediaapi.Wikipedia(language=lang,extract_format=wikipediaapi.ExtractFormat.WIKI, data={"action":"parse"}) # Create wikipedia class
        self.pages = self.parse_raw_topics(topics) # Parse topics
        self.presentation = pptx.Presentation() # Create presentation class

        self.generate_content(self.pages, self.presentation) # Generate content
        self.save_content() # Save content

    def parse_raw_topics(self, topics):
        result = []
        for topic in topics: 
            if "#" in topic: 
                topic = topic.split("#")
                page = self.wikipedia.page(topic[0])
                section = page.section_by_title(topic[1])
                if page.exists():
                    result.append(section)
                    if section == None:
                        raise Exception("Couldn't find a section")
                    self.contents.append(topic[1])
                else:
                    raise Exception("Couldn't find a page")
            else:
                page = self.wikipedia.page(topic)
                if page.exists():
                    titles = self.recursively_find_sections(page.sections)
                    titles = flatten(titles)
                    for title in titles:
                        result.append(page.section_by_title(title))
                        self.contents.append(title)
                else:
                    raise Exception("Couldn't find a page")

        return result
    def recursively_find_sections(self, sections):
        result=[]
        for s in sections:
            result.append(s.title)
            rec = self.recursively_find_sections(s.sections)
            if rec != []:
                result.append(rec)
        return result

    def generate_content(self, pages, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.placeholders[0].text = self.title
        slide.placeholders[1].text = self.author
    
        contentslide=prs.slides.add_slide(prs.slide_layouts[1])
        contentslide.placeholders[0].text = "Spis treści"
        contentslide.placeholders[1].text = "\n".join(self.contents)

        for page in pages:
            count = 1
            last_count = 0
            for c in page.text: ## replace this with a custom regex, for example \(\.) ([A-Z]\w+)\g (searching for end of sentence, "Goodbye, my friend[. N]ext sentence")
                if c == self.split_on and last_count > self.max_letters_per_page:
                    count+=1;
                    last_count = 0
                last_count+=1
            amount_of_slices = count
            for i in range(amount_of_slices):
                lst = [0]
                last_count = 0
                for pos,char in enumerate(page.text):
                    if char == self.split_on and last_count > self.max_letters_per_page:
                        lst.append(pos)
                        last_count = 0
                    last_count+=1
                lst.append(len(page.text))

                if len(lst) != 0:
                    x = lst[i]
                    if len(lst) != 1:
                        y = lst[i+1]
                    else:
                        y = len(page.text)
                    ptext = page.text[slice(x,y)].replace("\n","").strip("\n").strip(self.split_on)
                    if ptext != "":
                        slide = prs.slides.add_slide(prs.slide_layouts[1])
                        slide.shapes.title.text = f"""{page.title} ({i+1}/{amount_of_slices})"""
                        slide.placeholders[1].text = ptext                 
                        slide.placeholders[1].text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                        
    def save_content(self):
        self.presentation.save(self.path)
